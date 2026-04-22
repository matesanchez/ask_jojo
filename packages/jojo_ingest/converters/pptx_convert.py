"""PPTX → markdown.

Each slide becomes an `## Slide N` section. Text runs are concatenated in
reading order; bullet lists are preserved; speaker notes are included under
a `**Speaker notes**:` subheading when present.

Images are skipped — same rationale as docx_convert. Tables use the same
markdown-table rendering as xlsx_convert.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text and para.text:
            text = para.text.strip()
        if not text:
            continue
        # Indent level → nested markdown list indentation.
        if para.level and para.level > 0:
            prefix = "    " * para.level + "- "
        else:
            prefix = "- " if _looks_bulleted(para) else ""
        lines.append(prefix + text)
    return "\n".join(lines)


def _looks_bulleted(para) -> bool:
    # python-pptx doesn't expose bullet formatting cleanly across all themes;
    # heuristic: any paragraph at level >= 0 in a text frame that has siblings.
    return para.level is not None


def _shape_table(shape) -> str | None:
    if not shape.has_table:
        return None
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text.replace("|", "\\|").replace("\n", " ") for cell in row.cells])
    if not rows:
        return None
    width = max(len(r) for r in rows)
    header = rows[0] + [""] * (width - len(rows[0]))
    body = rows[1:]
    lines = ["| " + " | ".join(header) + " |",
             "| " + " | ".join(["---"] * width) + " |"]
    for r in body:
        padded = r + [""] * (width - len(r))
        lines.append("| " + " | ".join(padded) + " |")
    return "\n".join(lines)


def convert_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    sections: list[str] = [f"# {path.stem}\n"]
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                table_md = _shape_table(shape)
                if table_md:
                    body_parts.append(table_md)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = _shape_text(shape).strip()
            if not text:
                continue
            if not title and shape == slide.shapes.title:
                title = text
            else:
                body_parts.append(text)
        header = f"## Slide {i}" + (f" — {title}" if title else "")
        section = [header]
        if body_parts:
            section.append("\n\n".join(body_parts))
        notes = _extract_notes(slide)
        if notes:
            section.append(f"\n**Speaker notes:** {notes}")
        sections.append("\n\n".join(section))
    return "\n\n".join(sections) + "\n"


def _extract_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    notes_tf = slide.notes_slide.notes_text_frame
    if not notes_tf:
        return ""
    text = notes_tf.text.strip()
    return text
