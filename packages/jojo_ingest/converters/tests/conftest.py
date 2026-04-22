"""Test fixtures that build tiny real Office/PDF files for converter tests.

We don't want to ship binary fixtures — keeps the repo lean and makes it
obvious what each converter is being asked to handle.
"""

from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    from docx import Document

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Sample Protocol", level=1)
    doc.add_paragraph(
        "This is the first paragraph. It describes the assay in plain prose."
    )
    doc.add_heading("Procedure", level=2)
    doc.add_paragraph("Step one — prepare the buffer.", style="List Bullet")
    doc.add_paragraph("Step two — add the substrate.", style="List Bullet")
    doc.save(str(path))
    return path


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    from openpyxl import Workbook

    path = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Readings"
    ws.append(["Time (min)", "OD600", "Notes"])
    ws.append([0, 0.02, "start"])
    ws.append([15, 0.11, ""])
    ws.append([30, 0.24, "log phase"])
    ws2 = wb.create_sheet("Calibration")
    ws2.append(["Standard", "Value"])
    ws2.append(["A", 1.0])
    wb.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = "Plasma pharmacokinetics overview"

    # Content slide
    bullet_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_layout)
    slide2.shapes.title.text = "Key metrics"
    tf = slide2.placeholders[1].text_frame
    tf.text = "Cmax — peak plasma concentration"
    p2 = tf.add_paragraph()
    p2.text = "Tmax — time to peak"
    p3 = tf.add_paragraph()
    p3.text = "AUC — area under the curve"
    prs.save(str(path))
    return path


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    import fitz  # pymupdf

    path = tmp_path / "sample.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(
        (72, 72),
        "This is a short PDF used to test text extraction.\n"
        "It has two sentences in one paragraph.",
        fontsize=11,
    )
    page2 = doc.new_page()
    page2.insert_text((72, 72), "A second page with its own content.", fontsize=11)
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def sample_text(tmp_path: Path) -> Path:
    path = tmp_path / "notes.txt"
    path.write_text("Plain-text note.\nSecond line.\n", encoding="utf-8")
    return path
