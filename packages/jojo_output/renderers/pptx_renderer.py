"""PPTX renderer -- non-Marp PowerPoint deck rendering.

For Marp-style decks (the common case for "make me slides"), use
``renderers/marp.py`` and the frontend Web Worker. This renderer is
for the rarer case where the output needs to be a real .pptx file
(e.g. exported for someone outside JoJo Bot to edit, or run through
the corporate template).

Wraps python-pptx behind a typed spec. The spec mirrors PowerPoint's
own slide-layout model (title, content, two-content, etc.).

Public API:

- ``PptxSlide`` / ``PptxSpec``
- ``render_pptx(spec, out_path)``
"""

from __future__ import annotations

from pathlib import Path
from typing import Literal

from pydantic import BaseModel, Field


class PptxSlide(BaseModel):
    """One slide in the deck.

    Attributes:
        layout: PowerPoint layout enum (``title`` for cover,
            ``title_content`` for the most common slide, ``two_content``
            for side-by-side, ``section`` for divider).
        title: slide title.
        body: free-form prose body (rendered into the content
            placeholder).
        bullets: bullets (rendered into the content placeholder; takes
            precedence over body).
        right_body / right_bullets: only for ``two_content`` layout.
        notes: speaker notes.
    """

    layout: Literal["title", "title_content", "two_content", "section"] = "title_content"
    title: str = ""
    body: str = ""
    bullets: list[str] = Field(default_factory=list)
    right_body: str = ""
    right_bullets: list[str] = Field(default_factory=list)
    notes: str = ""


class PptxSpec(BaseModel):
    """Deck spec.

    Attributes:
        title: deck title (used for the cover slide).
        author: shown on the cover.
        slides: ordered slides.
    """

    title: str = ""
    author: str = "JoJo Bot"
    slides: list[PptxSlide] = Field(default_factory=list)


def _set_text(placeholder, text: str, bullets: list[str]) -> None:
    """Fill a placeholder with body text or bullets."""
    if not (text or bullets):
        return
    tf = placeholder.text_frame
    if bullets:
        # First bullet replaces the default empty paragraph.
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
    else:
        tf.text = text


def render_pptx(spec: PptxSpec, out_path: Path | str) -> Path:
    """Render a .pptx file. Requires ``python-pptx``."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is required for pptx rendering; install via "
            "pip install -e \".[output]\""
        ) from e

    prs = Presentation()

    # Cover slide if a title was provided.
    if spec.title:
        layout = prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = spec.title
        if len(s.placeholders) > 1:
            s.placeholders[1].text = spec.author

    layout_map = {
        "title": 0,
        "title_content": 1,
        "two_content": 3,
        "section": 2,
    }

    for slide_spec in spec.slides:
        layout_idx = layout_map.get(slide_spec.layout, 1)
        # python-pptx can crash on some templates if the layout index is
        # out of range; clamp.
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = slide_spec.title

        if slide_spec.layout == "two_content":
            # placeholders[1] = left, placeholders[2] = right (template-dependent)
            if len(slide.placeholders) >= 2:
                _set_text(slide.placeholders[1], slide_spec.body, slide_spec.bullets)
            if len(slide.placeholders) >= 3:
                _set_text(slide.placeholders[2], slide_spec.right_body, slide_spec.right_bullets)
        elif slide_spec.layout in ("title_content", "section"):
            if len(slide.placeholders) >= 2:
                _set_text(slide.placeholders[1], slide_spec.body, slide_spec.bullets)

        if slide_spec.notes and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = slide_spec.notes

    out = Path(out_path).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
